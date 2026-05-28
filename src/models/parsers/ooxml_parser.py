import os
import typing
try:
    import pptx
except ImportError:
    pptx = None
try:
    import docx
except ImportError:
    docx = None

def parse_ooxml_file(file_path: str) -> typing.Dict[str, int]:
    """
    FR-1.1, FR-1.2: .pptx / .docx / .hwpx 파일 파싱 및 글자 수 추출. 
    {"팀원이름": 총_글자수} 반환. 식별자 비어있으면 "Unknown".
    손상된 파일은 skip + 빈 dict 반환 (FR-1.1 수용기준).
    """
    if not os.path.exists(file_path):
        return {}
        
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".docx":
        return _parse_docx(file_path)
    elif ext == ".pptx":
        return _parse_pptx(file_path)
    elif ext == ".hwpx":
        return _parse_hwpx(file_path)
        
    return {}

def _parse_docx(file_path: str) -> typing.Dict[str, int]:
    """FR-1.1: .docx 파싱"""
    if not docx: return {}
    try:
        doc = docx.Document(file_path)
        author = doc.core_properties.author
        if not author:
            author = "Unknown"
            
        total_chars = 0
        for para in doc.paragraphs:
            text = para.text
            clean_text = text.replace(' ', '').replace('\n', '')
            total_chars += len(clean_text)
            
        return {author: total_chars}
    except Exception:
        return {"Unknown": 0}

def _parse_pptx(file_path: str) -> typing.Dict[str, int]:
    """FR-1.1: .pptx 파싱"""
    if not pptx: return {}
    try:
        prs = pptx.Presentation(file_path)
        author = prs.core_properties.last_modified_by
        if not author:
            author = "Unknown"
            
        total_chars = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text
                    clean_text = text.replace(' ', '').replace('\n', '')
                    total_chars += len(clean_text)
                    
        return {author: total_chars}
    except Exception:
        return {"Unknown": 0}

def _parse_hwpx(file_path: str) -> typing.Dict[str, int]:
    """
    FR-1.1: .hwpx 파싱 (HWPX는 OOXML 기반 ZIP 아카이브)
    HWPX 표준 스키마에서 텍스트를 추출한다.
    """
    import zipfile
    try:
        from xml.etree import ElementTree as ET
    except ImportError:
        return {"Unknown": 0}
        
    try:
        total_chars = 0
        author = "Unknown"
        
        with zipfile.ZipFile(file_path, 'r') as zf:
            # HWPX 메타데이터에서 작성자 추출 시도
            meta_files = [n for n in zf.namelist() if 'meta' in n.lower() or 'core' in n.lower()]
            for mf in meta_files:
                try:
                    meta_xml = zf.read(mf).decode('utf-8')
                    meta_root = ET.fromstring(meta_xml)
                    # dc:creator 또는 creator 태그 탐색
                    for elem in meta_root.iter():
                        tag_local = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                        if tag_local.lower() in ('creator', 'author') and elem.text:
                            author = elem.text.strip()
                            break
                    if author != "Unknown":
                        break
                except Exception:
                    continue
            
            # HWPX 본문 텍스트 추출: Contents/ 하위 section XML 파일들
            content_files = [n for n in zf.namelist() 
                           if n.startswith('Contents/') and n.endswith('.xml')]
            if not content_files:
                # 대안: 모든 XML 파일에서 텍스트 탐색
                content_files = [n for n in zf.namelist() if n.endswith('.xml')]
                
            for cf in content_files:
                try:
                    xml_data = zf.read(cf).decode('utf-8')
                    root = ET.fromstring(xml_data)
                    # 모든 텍스트 노드에서 문자 추출
                    for elem in root.iter():
                        if elem.text:
                            clean = elem.text.replace(' ', '').replace('\n', '').replace('\t', '').replace('\r', '')
                            total_chars += len(clean)
                except Exception:
                    continue
                    
        return {author: total_chars}
    except Exception:
        # 손상된 파일 → skip, 빈 결과 반환
        return {"Unknown": 0}
