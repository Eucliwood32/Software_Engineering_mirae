"""FR-1.1/1.2 DocumentParser: .docx/.pptx/.hwpx 텍스트·작성자 추출."""
from __future__ import annotations
import re
import zipfile


def _strip_whitespace(text: str) -> str:
    return re.sub(r"\s", "", text)


class DocumentParser:
    """FR-1.1, FR-1.2 — 확장자별 추출기로 위임하는 파사드. 손상 파일 → {} 반환."""

    def parse(self, path: str) -> dict[str, int]:
        """확장자에 따라 추출기로 위임. 반환 {작성자: 유효글자수}."""
        ext = path.lower().rsplit(".", 1)[-1]
        try:
            if ext == "docx":
                return self._parse_docx(path)
            elif ext == "pptx":
                return self._parse_pptx(path)
            elif ext == "hwpx":
                return self._parse_hwpx(path)
            return {}
        except Exception:
            return {}

    def count_shapes(self, path: str) -> int:
        """텍스트박스 포함 도형 개수 (.pptx 한정, 그 외 0)."""
        ext = path.lower().rsplit(".", 1)[-1]
        try:
            if ext == "pptx":
                return self._count_pptx_shapes(path)
            return 0
        except Exception:
            return 0

    # ── docx ──────────────────────────────────────────────────────────────

    def _parse_docx(self, path: str) -> dict[str, int]:
        from docx import Document
        doc = Document(path)
        author = doc.core_properties.author or "Unknown"
        text = " ".join(p.text for p in doc.paragraphs)
        return {author: len(_strip_whitespace(text))}

    # ── pptx ──────────────────────────────────────────────────────────────

    def _parse_pptx(self, path: str) -> dict[str, int]:
        from pptx import Presentation
        prs = Presentation(path)
        author = prs.core_properties.last_modified_by or "Unknown"
        total = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            total += len(_strip_whitespace(run.text))
        return {author: total}

    def _count_pptx_shapes(self, path: str) -> int:
        from pptx import Presentation
        prs = Presentation(path)
        return sum(
            1
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )

    # ── hwpx ──────────────────────────────────────────────────────────────

    def _parse_hwpx(self, path: str) -> dict[str, int]:
        """HWPX(OWPML zip)에서 작성자·본문 텍스트 추출 (FR-1.1, D-4)."""
        import xml.etree.ElementTree as ET

        with zipfile.ZipFile(path, "r") as zf:
            author = "Unknown"
            for name in zf.namelist():
                if "metadata" in name.lower() or "meta-inf" in name.lower():
                    try:
                        root = ET.fromstring(zf.read(name))
                        for elem in root.iter():
                            tag = elem.tag.split("}")[-1].lower()
                            if tag == "creator" and elem.text:
                                author = elem.text.strip() or "Unknown"
                                break
                    except Exception:
                        pass
                if author != "Unknown":
                    break

            total = 0
            for name in zf.namelist():
                if re.search(r"section\d*\.xml$", name, re.IGNORECASE):
                    try:
                        root = ET.fromstring(zf.read(name))
                        for elem in root.iter():
                            tag = elem.tag.split("}")[-1]
                            if tag == "t" and elem.text:
                                total += len(_strip_whitespace(elem.text))
                    except Exception:
                        pass

        return {author: total}
