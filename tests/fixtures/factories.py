"""
테스트 픽스처 팩토리 (test-plan §4).
구현 완료 전까지 각 함수는 NotImplementedError를 발생시킨다.
팩토리가 GREEN이 된 뒤 L1 단위 테스트를 시작할 수 있다.
"""
from __future__ import annotations

import zipfile
import os


# ---------------------------------------------------------------------------
# FR-1  문서 팩토리
# ---------------------------------------------------------------------------

def make_docx(path: str, author: str | None, text: str) -> str:
    """author=None이면 core_properties.author를 비운다(Unknown 검증용).
    text는 단락 1개로 기록. 반환: path."""
    from docx import Document
    doc = Document()
    doc.add_paragraph(text)
    doc.core_properties.author = author if author is not None else ""
    doc.save(path)
    return path


def make_pptx(path: str, last_modified_by: str | None,
              slides: list[list[str]]) -> str:
    """slides = [[box1_text, box2_text, ...], ...]. 슬라이드당 텍스트박스 N개 생성.
    last_modified_by=None이면 메타 공란. 반환: path."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for boxes in slides:
        slide = prs.slides.add_slide(blank)
        for i, t in enumerate(boxes):
            tb = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(4), Inches(1))
            tb.text_frame.text = t
    if last_modified_by is not None:
        prs.core_properties.last_modified_by = last_modified_by
    prs.save(path)
    return path


def make_hwpx(path: str, author: str | None, text: str) -> str:
    """최소 HWPX(zip) 생성. 표준 메타 스키마에 author 기록, 본문 text 1문단.
    반환: path. (HWPX = OWPML zip; mimetype + Contents/ + 메타 XML)"""
    # OWPML 최소 구조
    mimetype = "application/hwp+zip"

    version_xml = '<?xml version="1.0" encoding="UTF-8"?>\n<hwpml version="1.0"/>'

    author_val = author if author is not None else ""
    meta_xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<opf:metadata xmlns:opf="http://www.idpf.org/2007/opf"
              xmlns:dc="http://purl.org/dc/elements/1.1/">
  <dc:creator>{author_val}</dc:creator>
</opf:metadata>"""

    # 본문 텍스트를 담는 최소 section XML
    section_xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<hml:sec xmlns:hml="http://www.hancom.co.kr/hwpml/2012/Hwp">
  <hml:p><hml:run><hml:t>{text}</hml:t></hml:run></hml:p>
</hml:sec>"""

    content_hpf = """<?xml version="1.0" encoding="UTF-8"?>
<package version="1.0" unique-identifier="hwpDocId"
         xmlns="http://www.idpf.org/2007/opf">
  <manifest>
    <item id="section0" href="section0.xml" media-type="application/xml"/>
    <item id="meta" href="../META-INF/metadata.xml" media-type="application/xml"/>
  </manifest>
</package>"""

    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", mimetype)
        zf.writestr("version.xml", version_xml)
        zf.writestr("META-INF/metadata.xml", meta_xml)
        zf.writestr("Contents/content.hpf", content_hpf)
        zf.writestr("Contents/section0.xml", section_xml)
    return path


def make_corrupted(path: str) -> str:
    """확장자만 .docx/.pptx/.hwpx, 내용은 무작위 바이트(zip 헤더 아님)."""
    with open(path, "wb") as f:
        f.write(b"\x00\x01NOT_A_ZIP\xff")
    return path


def make_empty_docx(path: str) -> str:
    """단락 0개짜리 빈 docx."""
    from docx import Document
    doc = Document()
    doc.save(path)
    return path


# ---------------------------------------------------------------------------
# FR-2  Git 저장소 팩토리
# ---------------------------------------------------------------------------

def make_git_repo(path: str, commits: list[dict]) -> str:
    """commits[i] = {"email": str, "date": "YYYY-MM-DD HH:MM:SS",
                     "add": int, "del": int}
    각 커밋마다 add줄 추가 / del줄 삭제를 정확히 발생시킨다.
    반환: path."""
    import subprocess

    subprocess.run(["git", "init", "-q", path], check=True)
    subprocess.run(["git", "-C", path, "config", "user.email", "setup@test.com"], check=True)
    subprocess.run(["git", "-C", path, "config", "user.name", "Setup"], check=True)

    # 베이스 파일 생성 (첫 커밋 이전에 충분한 줄 수를 확보하기 위해)
    base_file = os.path.join(path, "work.txt")

    # 모든 커밋에서 필요한 최대 del 합산 수만큼 베이스 라인 확보
    max_needed = sum(c.get("del", 0) for c in commits) + 100
    with open(base_file, "w", encoding="utf-8") as f:
        for i in range(max_needed):
            f.write(f"base_line_{i}\n")

    env_base = {
        "GIT_AUTHOR_NAME": "Base",
        "GIT_AUTHOR_EMAIL": "base@setup.internal",
        "GIT_AUTHOR_DATE": "2000-01-01T00:00:00",
        "GIT_COMMITTER_NAME": "Base",
        "GIT_COMMITTER_EMAIL": "base@setup.internal",
        "GIT_COMMITTER_DATE": "2000-01-01T00:00:00",
        **os.environ,
    }
    subprocess.run(["git", "-C", path, "add", "."], check=True)
    subprocess.run(["git", "-C", path, "commit", "-m", "base", "--allow-empty"],
                   check=True, env=env_base)

    current_line_count = max_needed

    for commit in commits:
        email: str = commit["email"]
        date: str = commit["date"]
        add_n: int = commit.get("add", 0)
        del_n: int = commit.get("del", 0)

        # del_n줄 삭제 후 add_n줄 추가
        with open(base_file, "r", encoding="utf-8") as f:
            lines = f.readlines()

        # 앞에서 del_n줄 제거
        lines = lines[del_n:]
        # add_n줄 추가
        new_lines = [f"new_line_{i}_{email}\n" for i in range(add_n)]
        lines = lines + new_lines

        with open(base_file, "w", encoding="utf-8") as f:
            f.writelines(lines)

        current_line_count = len(lines)

        name = email.split("@")[0]
        iso_date = date.replace(" ", "T")
        env = {
            "GIT_AUTHOR_NAME": name,
            "GIT_AUTHOR_EMAIL": email,
            "GIT_AUTHOR_DATE": iso_date,
            "GIT_COMMITTER_NAME": name,
            "GIT_COMMITTER_EMAIL": email,
            "GIT_COMMITTER_DATE": iso_date,
            **os.environ,
        }
        subprocess.run(["git", "-C", path, "add", "."], check=True, env=env)
        subprocess.run(
            ["git", "-C", path, "commit", "-m", f"commit by {email}"],
            check=True, env=env,
        )

    return path


# ---------------------------------------------------------------------------
# FR-3  카카오톡 텍스트 팩토리
# ---------------------------------------------------------------------------

def make_katalk(path: str,
                lines: list[tuple[str, str]] | list[str],
                encoding: str = "utf-8",
                header_date: str = "2024년 1월 15일 월요일") -> str:
    """lines가 (author, msg) 튜플이면 정상 발화줄로 직렬화.
    lines가 str이면 그 줄을 그대로(오염줄 주입용) 기록.
    encoding='cp949' 로 CP949 파일 생성 가능."""
    output_lines = [header_date]
    for item in lines:
        if isinstance(item, tuple):
            author, msg = item
            output_lines.append(f"[{author}] [오후 2:30] {msg}")
        else:
            output_lines.append(item)

    content = "\n".join(output_lines) + "\n"
    with open(path, "w", encoding=encoding) as f:
        f.write(content)
    return path


# ---------------------------------------------------------------------------
# FR-5  MemberScore 샘플 픽스처
# ---------------------------------------------------------------------------

def sample_scores(n: int = 4):
    """차트 테스트용 MemberScore 리스트. 구현 전이면 ImportError로 RED."""
    from qce.model.types import MemberScore
    return [
        MemberScore("조원희", 0.30, 0.20, 0.40, 0.28, 800, 1500, 120, False, []),
        MemberScore("B팀원",  0.55, 0.10, 0.30, 0.31, 1000, 400, 90, True, ["EW-01"]),
        MemberScore("C팀원",  0.40, 0.45, 0.20, 0.29, 600, 2000, 50, False, []),
        MemberScore("D팀원",  0.05, 0.08, 0.05, 0.12, 50, 200, 10, False, ["ZSCORE"]),
    ][:n]
