"""FR-1.1/1.2 DocumentParser: .docx/.pptx/.hwpx 텍스트·작성자 추출."""
from __future__ import annotations
import zipfile
import re


def _strip_whitespace(text: str) -> str:
    """공백·탭·개행 제외 유효 글자만 카운트."""
    return re.sub(r"\s", "", text)


class DocumentParser:
    def parse(self, path: str) -> dict[str, int]:
        """확장자별 추출기로 위임. {작성자: 유효글자수}. 손상 → {} 반환."""
        ext = path.lower().rsplit(".", 1)[-1]
        try:
            if ext == "docx":
                return self._parse_docx(path)
            elif ext == "pptx":
                return self._parse_pptx(path)
            elif ext == "hwpx":
                return self._parse_hwpx(path)
            else:
                return {}
        except Exception:
            return {}

    def count_shapes(self, path: str) -> int:
        """문서 내 텍스트박스 포함 도형 개수."""
        ext = path.lower().rsplit(".", 1)[-1]
        try:
            if ext == "pptx":
                return self._count_pptx_shapes(path)
            return 0
        except Exception:
            return 0

    # ── docx ──────────────────────────────────────────────────────────────

    def _parse_docx(self, path: str) -> dict[str, int]:
        from docx import Document
        doc = Document(path)
        author = doc.core_properties.author or "Unknown"
        text = " ".join(p.text for p in doc.paragraphs)
        return {author: len(_strip_whitespace(text))}

    # ── pptx ──────────────────────────────────────────────────────────────

    def _parse_pptx(self, path: str) -> dict[str, int]:
        from pptx import Presentation
        prs = Presentation(path)
        author = prs.core_properties.last_modified_by or "Unknown"
        total = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            total += len(_strip_whitespace(run.text))
        return {author: total}

    def _count_pptx_shapes(self, path: str) -> int:
        from pptx import Presentation
        prs = Presentation(path)
        count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    count += 1
        return count

    # ── hwpx ──────────────────────────────────────────────────────────────

    def _parse_hwpx(self, path: str) -> dict[str, int]:
        """HWPX(OWPML zip)에서 작성자·본문 텍스트 추출."""
        import xml.etree.ElementTree as ET

        with zipfile.ZipFile(path, "r") as zf:
            # 작성자: META-INF/metadata.xml의 dc:creator
            author = "Unknown"
            meta_candidates = [
                n for n in zf.namelist()
                if "metadata" in n.lower() or "meta-inf" in n.lower()
            ]
            for mc in meta_candidates:
                try:
                    xml_bytes = zf.read(mc)
                    root = ET.fromstring(xml_bytes)
                    for elem in root.iter():
                        tag = elem.tag.split("}")[-1].lower()
                        if tag == "creator" and elem.text:
                            author = elem.text.strip() or "Unknown"
                            break
                except Exception:
                    pass

            # 본문: Contents/section*.xml의 <hml:t> 텍스트
            total = 0
            section_files = [
                n for n in zf.namelist()
                if re.search(r"section\d*\.xml$", n, re.IGNORECASE)
            ]
            for sf in section_files:
                try:
                    xml_bytes = zf.read(sf)
                    root = ET.fromstring(xml_bytes)
                    for elem in root.iter():
                        tag = elem.tag.split("}")[-1]
                        if tag == "t" and elem.text:
                            total += len(_strip_whitespace(elem.text))
                except Exception:
                    pass

        return {author: total}
